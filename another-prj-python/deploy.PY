from pptx import Presentation

# Create a presentation object
presentation = Presentation()

# Add a title slide
slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])  # Title slide layout
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Sample Presentation"
subtitle.text = "Created using Python"

# Add a content slide
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])  # Content slide layout
title = slide2.shapes.title
content = slide2.placeholders[1]

title.text = "Agenda"
content.text = "1. Introduction\n2. Python-Powered Slides\n3. Summary"

# Add a slide with a picture (optional)
slide3 = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank layout
title = slide3.shapes.title
img_path = "example_image.png"  # Provide your image path
try:
    slide3.shapes.add_picture(img_path, left=0, top=0, width=presentation.slide_width, height=None)
except FileNotFoundError:
    print(f"Image not found at {img_path}. Skipping this slide.")

# Save the presentation
presentation.save("output_presentation.pptx")
print("Presentation saved as 'output_presentation.pptx'")
